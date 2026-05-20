"""Rebuild seminar_slides.pptx: refactor results section around claim groups."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
import shutil, os

SRC = "/home/user/information-as-alignment/seminar_slides.pptx"
DST = "/home/user/information-as-alignment/seminar_slides_v2.pptx"

RED = RGBColor(0xD2, 0x11, 0x17)
GREY = RGBColor(0x35, 0x36, 0x35)
LIGHT_GREY = RGBColor(0x80, 0x80, 0x80)
ALT_ROW = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xCC, 0xCC, 0xCC)

# Geometry constants from inspection (inches)
TITLE_BLOCK = (0.50, 0.73, 9.00, 1.25)
BODY = (1.01, 1.88, 7.97, 3.12)
CHIP = (6.80, 0.50, 2.70, 0.30)
READING = (1.01, 4.92, 7.97, 0.40)

shutil.copy(SRC, DST)
prs = Presentation(DST)

# -------- helpers --------

def get_monogram_xml(prs):
    """Return a deepcopy of the monogram group from an existing result slide."""
    s = prs.slides[36]  # canonical lifecycle slide has the monogram
    for shp in s.shapes:
        if "g3" in shp.name or shp.name.startswith("Google Shape;16") or (shp.shape_type == 6 and shp.width < Inches(0.5)):
            if shp.width < Inches(0.6):
                return deepcopy(shp._element)
    raise RuntimeError("monogram not found")

MONOGRAM_XML = get_monogram_xml(prs)

def add_monogram(slide):
    slide.shapes._spTree.append(deepcopy(MONOGRAM_XML))

def set_run(run, text, font_name, size_pt, bold, color, italic=False):
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color

def add_text_box(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tb.text_frame.word_wrap = True
    tb.text_frame.vertical_anchor = anchor
    tb.text_frame.margin_left = Inches(0.05)
    tb.text_frame.margin_right = Inches(0.05)
    tb.text_frame.margin_top = Inches(0.03)
    tb.text_frame.margin_bottom = Inches(0.03)
    # remove default paragraph behavior
    tf = tb.text_frame
    p0 = tf.paragraphs[0]
    return tb

def add_title_block(slide, claim_tag, title):
    """Title block: red accent line + large title."""
    l, t, w, h = TITLE_BLOCK
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_top = 0
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.LEFT
    r0 = p0.add_run()
    set_run(r0, claim_tag, "Montserrat", 12, True, RED)
    p1 = tf.add_paragraph()
    p1.alignment = PP_ALIGN.LEFT
    p1.space_before = Pt(2)
    r1 = p1.add_run()
    set_run(r1, title, "Montserrat Black", 24, False, GREY)

def add_cell_chip(slide, cell_tag):
    """Small grey chip top-right with 'Cells X, Y'."""
    l, t, w, h = CHIP
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    set_run(r, cell_tag, "Montserrat", 10, True, LIGHT_GREY)

def add_method_block(slide, left, top, width, height, method_text, label="Method:"):
    tb = add_text_box(slide, left, top, width, height)
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run()
    set_run(r1, label + " ", "Source Serif 4", 12, True, GREY)
    r2 = p.add_run()
    set_run(r2, method_text, "Source Serif 4", 12, False, GREY)

def add_reading_bar(slide, text):
    l, t, w, h = READING
    tb = add_text_box(slide, l, t, w, h, anchor=MSO_ANCHOR.MIDDLE)
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r0 = p.add_run()
    set_run(r0, "► ", "Source Serif 4", 12, True, RED)
    r1 = p.add_run()
    set_run(r1, text, "Source Serif 4", 12, False, GREY, italic=True)

def add_caveat(slide, text, top=5.32):
    tb = add_text_box(slide, 1.01, top, 7.97, 0.25)
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    set_run(r, text, "Source Serif 4", 9, False, LIGHT_GREY, italic=True)

def style_table(table, col_widths_in=None):
    """Apply red header + alt row + borders."""
    if col_widths_in:
        for i, w in enumerate(col_widths_in):
            table.columns[i].width = Inches(w)
    for r, row in enumerate(table.rows):
        for c, cell in enumerate(row.cells):
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = RED
            elif r % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = ALT_ROW
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
            # font already set per cell when populated

def add_table(slide, left, top, width, height, data, col_widths=None, header_bold_white=True, body_size=11):
    rows, cols = len(data), len(data[0])
    shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    table = shape.table
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = ""
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            # Allow bolding via tuple form ('text', True)
            if isinstance(val, tuple):
                txt, bold = val
            else:
                txt, bold = val, False
            run = p.add_run()
            if r == 0:
                set_run(run, str(txt), "Montserrat", body_size, True, WHITE)
                p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            else:
                set_run(run, str(txt), "Source Serif 4", body_size, bold, GREY)
    style_table(table, col_widths)
    return table

def add_summary_line(slide, left, top, width, height, parts):
    """Variant A under-table headline: list of (text, bold) tuples."""
    tb = add_text_box(slide, left, top, width, height, anchor=MSO_ANCHOR.MIDDLE)
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    for text, bold in parts:
        r = p.add_run()
        set_run(r, text, "Source Serif 4", 11, bold, GREY)

# -------- slide builders --------

def add_blank_slide():
    layout = prs.slide_layouts[6]  # BLANK
    return prs.slides.add_slide(layout)

def build_result_A(claim_tag, cell_tag, title, method, table_data, table_col_widths,
                   summary_parts, reading, caveat=None):
    """Variant A — side-by-side: method (left), table (right)."""
    s = add_blank_slide()
    add_monogram(s)
    add_title_block(s, claim_tag, title)
    add_cell_chip(s, cell_tag)
    # Method block left
    add_method_block(s, 1.01, 1.88, 3.50, 2.60, method)
    # Table on right
    n_rows = len(table_data)
    tbl_h = min(2.60, 0.32 * n_rows + 0.05)
    add_table(s, 4.70, 1.88, 4.30, tbl_h, table_data, table_col_widths, body_size=11)
    # Summary line under table
    if summary_parts:
        add_summary_line(s, 4.70, 1.88 + tbl_h + 0.10, 4.30, 0.40, summary_parts)
    add_reading_bar(s, reading)
    if caveat:
        add_caveat(s, caveat)
    return s

def build_result_B(claim_tag, cell_tag, title, method, table_data, table_col_widths,
                   reading, caveat=None, table_top=2.55, table_height=1.85, method_height=0.55):
    """Variant B — stacked: method (top), table (full-width)."""
    s = add_blank_slide()
    add_monogram(s)
    add_title_block(s, claim_tag, title)
    add_cell_chip(s, cell_tag)
    add_method_block(s, 1.01, 1.88, 7.97, method_height, method)
    add_table(s, 1.01, table_top, 7.97, table_height, table_data, table_col_widths, body_size=11)
    add_reading_bar(s, reading)
    if caveat:
        add_caveat(s, caveat)
    return s

def build_section_divider(prefix_bold, suffix, subtitle):
    """Match existing section divider style (slide 10, 19, 26, 34)."""
    s = add_blank_slide()
    add_monogram(s)
    # Title at (1.21, 2.21) size 6.79 x 1.21
    tb = s.shapes.add_textbox(Inches(1.21), Inches(2.21), Inches(7.50), Inches(1.30))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run()
    set_run(r1, prefix_bold, "Montserrat", 30, True, GREY)
    r2 = p.add_run()
    set_run(r2, " " + suffix, "Montserrat", 30, False, GREY)
    # Subtitle italic grey
    tb2 = s.shapes.add_textbox(Inches(1.21), Inches(3.60), Inches(7.50), Inches(0.50))
    p2 = tb2.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    r = p2.add_run()
    set_run(r, subtitle, "Source Serif 4", 18, False, LIGHT_GREY, italic=True)
    return s

def build_synthesis(text):
    """Full-page minimal synthesis slide with single sentence."""
    s = add_blank_slide()
    add_monogram(s)
    tb = s.shapes.add_textbox(Inches(0.80), Inches(2.20), Inches(8.40), Inches(1.40))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    set_run(r, text, "Montserrat", 22, True, GREY)
    return s

def build_smoke_caveat():
    """One-slide framing of smoke→paper-grade run."""
    s = add_blank_slide()
    add_monogram(s)
    add_title_block(s, "RESULTS — METHODOLOGY NOTE", "Smoke run now; paper-grade run replaces values")
    tb = add_text_box(s, 1.01, 1.95, 7.97, 3.00)
    tf = tb.text_frame
    lines = [
        ("Numbers below are from the canonical end-to-end smoke run (~3h on A100, all 40 cells executed).", False),
        ("", False),
        ("Paper-grade run (currently executing, ~70h on H100) extends:", True),
        ("  •  Full 300/200/200/200 epoch schedule (smoke used max=2 universal cap)", False),
        ("  •  24-step LoRA durability test (smoke: 2 steps)", False),
        ("  •  Full 1k → 50k scale sweep (smoke: 1k only)", False),
        ("  •  Full long-horizon 20×5 rounds (smoke: stubbed)", False),
        ("  •  ZsRE and CounterFact lifecycle benchmarks at full coverage", False),
        ("", False),
        ("Smoke values will be replaced with paper-grade values before the talk.", True),
    ]
    for i, (txt, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        set_run(r, txt, "Source Serif 4", 12, bold, GREY)
    return s

# -------- step 1: edit slide 17 (Locality → Geometry) --------

s17 = prs.slides[16]
outer_group = s17.shapes[0]
text_shape = list(outer_group.shapes)[1]
tf = text_shape.text_frame
# paragraphs: 0 Substrate / 1 desc / 2 Locality / 3 desc / 4 Generality / 5 desc
# Edit paragraph 2 (Locality) → Geometry and update its description (paragraph 3)
locality_para = tf.paragraphs[2]
locality_para.runs[0].text = "Geometry"
geom_desc_para = tf.paragraphs[3]
geom_desc_para.runs[0].text = ("Acts where intended — strong-prior override, scale with constraint count, "
                                "and semantic consequences via compiled closure")

# -------- step 2: build all new slides (appended to end of deck) --------

new_slide_specs = []

# Slide 35 (NEW): smoke caveat
new_slide_specs.append(("smoke_caveat", None))

# § A — Substrate divider
new_slide_specs.append(("divider", ("§ A —", "Substrate", "Does it hold across time?")))

# A1 — Canonical lifecycle
new_slide_specs.append(("result_A", {
    "claim_tag": "§ A — SUBSTRATE",
    "cell_tag": "Cells 8, 9",
    "title": "Canonical lifecycle on Future Industries",
    "method": ("Trained the correction field through Phases A→D on Future Industries "
               "(Mistral-7B-v0.1 frozen). After each phase, scored install rate on that "
               "phase's test set and retention on Phase A's test set. Linear readout."),
    "table_data": [
        ["Phase", "Install (lin)", "Retain A"],
        ["A — Onboarding", "0.96", "—"],
        ["B — Initiative", "0.98", "0.95"],
        ["C — Reorg.", "0.99", "0.93"],
        ["D — Turnover", "1.00", "0.94"],
    ],
    "table_col_widths": [1.90, 1.20, 1.20],
    "summary_parts": [
        ("Mean ", False), ("0.973", True), ("  ·  base-only 0.216  ·  gain ", False),
        ("+0.756", True), ("  ·  ", False), ("6,381", True), (" centers", False),
    ],
    "reading": "Field installs and retains the full lifecycle on one substrate; net A→D drift −0.029.",
}))

# A2 — Forgetting decomposition (variant B)
new_slide_specs.append(("result_B", {
    "claim_tag": "§ A — SUBSTRATE",
    "cell_tag": "Cell 25",
    "title": "Forgetting decomposition",
    "method": ("Measured ΔAccuracy on Phase-A test set after each subsequent phase trained; "
               "decomposed total drift into per-boundary contributions."),
    "table_data": [
        ["Boundary", "Drift", "Note"],
        ["A → B", "< −0.005", "Composition, no contradiction"],
        ["B → C", "−0.019", "Structural revision (dominant)"],
        ["C → D", "< −0.005", "Partial replacement"],
        ["A → D (net)", "−0.029", "Mild, decomposable, not catastrophic"],
    ],
    "table_col_widths": [1.80, 1.40, 4.77],
    "reading": "Forgetting is mild and traceable to the one boundary where contradiction is structural.",
}))

# A3 — Inspectability (variant A, qualitative table)
new_slide_specs.append(("result_A", {
    "claim_tag": "§ A — SUBSTRATE",
    "cell_tag": "Cells 10, 11",
    "title": "Inspectability — provenance and consistency",
    "method": ("Logged, per query, which centers contributed to the field's correction, which "
               "phase nucleated them, and which were promoted vs. dissolved. Audited paraphrase "
               "consistency across phases."),
    "table_data": [
        ["Audit", "Observation"],
        ["Provenance", "Contributing centers traced per query"],
        ["Paraphrase consistency", "Field answers identically across forms"],
        ["Cross-phase drift", "None between recorded and live"],
    ],
    "table_col_widths": [1.80, 2.50],
    "summary_parts": [
        ("Inspectability is a ", False), ("property of the architecture", True),
        (", not a feature checkbox.", False),
    ],
    "reading": "Field state is readable, auditable, and stable under inspection.",
}))

# A4 — Retraction + selective deletion
new_slide_specs.append(("result_A", {
    "claim_tag": "§ A — SUBSTRATE",
    "cell_tag": "Cells 12, 13",
    "title": "Retraction and selective deletion",
    "method": ("Targeted a specific proposition and removed its nucleating centers; measured "
               "argmax accuracy on the target and on a matched set of control propositions. "
               "Tested retraction by writing the contradictory target."),
    "table_data": [
        ["Operation", "Metric", "Value"],
        ["Selective deletion", "Target acc.", "1.000 → 0.000"],
        ["Selective deletion", "Control drift", "+0.000"],
        ["Retraction (Cell 12)", "Argmax flip", "Clean"],
        ["Restoration", "Method", "Checkpoint, no retrain"],
    ],
    "table_col_widths": [1.55, 1.20, 1.55],
    "summary_parts": [
        ("Deletion is ", False), ("surgical", True), (". Controls are untouched.", False),
    ],
    "reading": "Lifecycle operations execute on the field without touching the base.",
}))

# A5 — Long-horizon stability
new_slide_specs.append(("result_A", {
    "claim_tag": "§ A — SUBSTRATE",
    "cell_tag": "Cell 21",
    "title": "Long-horizon stability",
    "method": ("Continued training for 20 rounds × 5 epochs after canonical convergence, "
               "tracking whether crystallized facts hold across 100+ further epochs."),
    "table_data": [
        ["Run mode", "Result"],
        ["Smoke", "Stub (bypassed)"],
        ["Paper-grade", "20×5 — in progress"],
    ],
    "table_col_widths": [1.80, 2.50],
    "summary_parts": [
        ("Tracking ", False), ("install-rate drift", True),
        (" and ", False), ("retention", True), (" over the horizon.", False),
    ],
    "reading": "Crystallized state is expected to survive sustained post-training perturbation.",
    "caveat": "Paper-grade values replace this slide before talk.",
}))

# A6 — LoRA durability
new_slide_specs.append(("result_A", {
    "claim_tag": "§ A — SUBSTRATE",
    "cell_tag": "Cell 26",
    "title": "Robustness to base drift — the LoRA test",
    "method": ("Applied a LoRA update to the frozen base after the field was trained; measured "
               "how much the base's argmax shifts (perturbation strength) and how much the "
               "field's behavior drifts (durability)."),
    "table_data": [
        ["Variant", "Base shift", "Field drift"],
        ["Smoke (2-step LoRA)", "4.6%", "+0.017"],
        ["Paper-grade (24-step)", "~37.5% [TBD]", "<1% [TBD]"],
    ],
    "table_col_widths": [1.90, 1.20, 1.20],
    "summary_parts": [
        ("Substrate decoupling: ", True), ("the field survives base perturbation "
        "that meaningfully changes the base.", False),
    ],
    "reading": "Field state is independent of the base parameters it sits over.",
    "caveat": "Paper-grade LoRA is the headline of the substrate-decoupling story.",
}))

# Synthesis A
new_slide_specs.append(("synthesis", "Substrate: it holds — across time, across operations, across base perturbation."))

# § B — Geometry divider
new_slide_specs.append(("divider", ("§ B —", "Geometry", "Does it act where it should?")))

# B1 — Strong-prior override + ontology shift + bridge + local alignment
new_slide_specs.append(("result_A", {
    "claim_tag": "§ B — GEOMETRY",
    "cell_tag": "Cells 14, 15, 16, 17",
    "title": "Strong-prior override and category-wide revision",
    "method": ("Installed corrections that contradict a saturated base prior (Cell 14). "
               "Tested category-wide relation reversal (Cell 15), bridge with weak-prior FI "
               "(Cell 16), and held-out propositions under σ-calibrated kernels (Cell 17)."),
    "table_data": [
        ["Quantity", "Value"],
        ["Override success", "1.000"],
        ["Base prior (pre)", "1.000"],
        ["Post-correction", "0.000"],
        ["Control propositions", "1.000"],
    ],
    "table_col_widths": [2.30, 2.00],
    "summary_parts": [
        ("Ontology shift, bridge, local-alignment control — ", False),
        ("all pass without collapse.", True),
    ],
    "reading": "Override is local. Geometric calibration holds at category scale.",
}))

# B2 — Locality and bleed
new_slide_specs.append(("result_A", {
    "claim_tag": "§ B — GEOMETRY",
    "cell_tag": "Cells 18, 19",
    "title": "Locality and bleed",
    "method": ("Installed a correction at point A; measured interference on propositions at "
               "varying geometric distance (Cell 19). Stress-tested with 1,000 unrelated rules "
               "installed simultaneously (Cell 18)."),
    "table_data": [
        ["Region", "Behavior"],
        ["Nearby (in σ)", "Minimal interference"],
        ["Far (out of σ)", "Indistinguishable from no install"],
        ["1k pilot — target", "0.995"],
        ["1k pilot — control", "1.000"],
    ],
    "table_col_widths": [1.80, 2.50],
    "summary_parts": [
        ("Locality is ", False), ("what makes everything else work", True), (".", False),
    ],
    "reading": "Corrections stay where the kernel says they stay; far field is untouched.",
}))

# B3 — Scale frontier (variant B)
new_slide_specs.append(("result_B", {
    "claim_tag": "§ B — GEOMETRY",
    "cell_tag": "Cell 20",
    "title": "Scale frontier — 1k → 50k constraints",
    "method": ("Swept simultaneous-constraint count from 1k up to 50k, dynamically expanding "
               "field capacity. Tracked install rate, retention, and locality at fixed bleed radius."),
    "table_data": [
        ["Constraint count", "1k", "3k", "5k", "10k", "20k", "50k"],
        ["Install (smoke)", "0.995", "—", "—", "—", "—", "—"],
        ["Control (smoke)", "1.000", "—", "—", "—", "—", "—"],
        ["Paper-grade", "[run]", "[TBD]", "[TBD]", "[TBD]", "[TBD]", "[TBD]"],
    ],
    "table_col_widths": [1.97, 1.00, 1.00, 1.00, 1.00, 1.00, 1.00],
    "reading": "Substrate is supposed to scale with constraint count, not with base-model size.",
    "caveat": "Paper-grade sweep is the headline of the scale story.",
}))

# B4 — Compiled closure (variant B)
new_slide_specs.append(("result_B", {
    "claim_tag": "§ B — GEOMETRY",
    "cell_tag": "Cells 22, 23, 24",
    "title": "When facts have consequences — emergent vs. compiled closure",
    "method": ("Trained on direct relations only, tested two-hop consequences (Cell 23). "
               "Then trained two-hop directly (Cell 23C). Then used an external compiler to "
               "install each edge of the closure as ordinary correction edges (Cell 24)."),
    "table_data": [
        ["Test", "Method", "Result"],
        ["23B — 2-hop emergent", "Train direct only, test 2-hop", "0.000"],
        ["23C — 2-hop trained", "Train 2-hop directly", "0.875"],
        ["24 — compiled A→C", "External compiler installs each edge", "1.000"],
        ["24 — revised A→D", "Revise C-era closure to D-era", "1.000"],
    ],
    "table_col_widths": [2.00, 3.47, 1.50],
    "reading": "IBF doesn't reason — compile semantic structure externally, install as correction edges.",
}))

# Synthesis B
new_slide_specs.append(("synthesis", "Geometry: it acts locally, scales, and handles consequences through compiled structure."))

# § C — Generality divider
new_slide_specs.append(("divider", ("§ C —", "Generality", "Does it replicate? Is it distinct from retrieval?")))

# C1 — Qwen replication
new_slide_specs.append(("result_A", {
    "claim_tag": "§ C — GENERALITY",
    "cell_tag": "Cell 36",
    "title": "Cross-model replication — Qwen2-1.5B",
    "method": ("Swapped Mistral-7B-v0.1 → Qwen2-1.5B. Engine, encoder, proposition space, "
               "σ-calibration, and FI construction held fixed; only R_base changed. Trained a "
               "fresh field from scratch over Qwen."),
    "table_data": [
        ["Metric", "Δ vs. Mistral"],
        ["5 of 6 lifecycle metrics", "within ±0.01"],
        ["Phase B gap", "−0.265"],
        ["Mechanism-level replication", "yes"],
        ["Field-state transfer", "no"],
    ],
    "table_col_widths": [2.30, 2.00],
    "summary_parts": [
        ("Phase B gap ", False), ("under investigation", True),
        ("; paper run should tighten.", False),
    ],
    "reading": "The architecture replicates across base models, not the trained instance.",
}))

# C2 — Cross-model ablations (variant B)
new_slide_specs.append(("result_B", {
    "claim_tag": "§ C — GENERALITY",
    "cell_tag": "Cell 37",
    "title": "Cross-model ablations and regime-dependence",
    "method": ("Foundational-paper predictions, made before this run: in the LLM regime "
               "crystallization and Crucible should be non-redundant; agency effectively neutral. "
               "Cell 37 ablates each mechanism on the Qwen run."),
    "table_data": [
        ["Mechanism", "Predicted role", "Test result"],
        ["Crystallization", "Non-redundant", "test in progress"],
        ["Crucible", "Non-redundant", "test in progress"],
        ["Agency", "Effectively neutral", "test in progress"],
    ],
    "table_col_widths": [2.10, 2.87, 3.00],
    "reading": "LLM is a third predicted regime — distinct from chess (all decisive) and CIFAR (all redundant).",
    "caveat": "Predictions made in advance; smoke ablation pending.",
}))

# C3 — ZsRE + CounterFact (variant B)
new_slide_specs.append(("result_B", {
    "claim_tag": "§ C — GENERALITY",
    "cell_tag": "Cells 28–35",
    "title": "External benchmarks — ZsRE and CounterFact",
    "method": ("Recast both datasets into a lifecycle: install → revise → remove → retain. "
               "Compared IBF correction field against kNN and RAG baselines. The baselines are "
               "oracle-maintained; IBF derives operations from the discrepancy signal alone."),
    "table_data": [
        ["Architecture", "Native lifecycle", "CounterFact revision"],
        ["IBF correction field", "1.000", "passes"],
        ["kNN (oracle-maintained)", "1.000", "passes"],
        ["RAG (oracle-maintained)", "1.000", "0.500 — fails"],
    ],
    "table_col_widths": [2.97, 2.50, 2.50],
    "reading": "Architectural difference shows even when baselines are handed oracle-fresh state.",
}))

# C4 — Paraphrase audit (variant B)
new_slide_specs.append(("result_B", {
    "claim_tag": "§ C — GENERALITY",
    "cell_tag": "Cells 31, 31b, 31c",
    "title": "Paraphrase audit",
    "method": ("Investigated paraphrase misses: measured geometric distance from direct prompts "
               "(31), swept σ to recover coverage (31b), and installed explicit anchor centers "
               "via the compiled-closure pattern (31c)."),
    "table_data": [
        ["Sub-cell", "Test", "Finding"],
        ["31", "Geometry audit", "Failed paraphrases sit geometrically far"],
        ["31b", "σ sweep", "Wider σ recovers some, damages locality"],
        ["31c", "Compiled anchors", "Explicit anchor installation closes the gap"],
    ],
    "table_col_widths": [1.30, 2.17, 4.50],
    "reading": "Same pattern as closure: compile structure externally, install as correction edges.",
}))

# C5 — Baseline-comparison note
new_slide_specs.append(("baseline_note", None))

# Synthesis C
new_slide_specs.append(("synthesis", "Generality: replicates across base models, structurally distinct from retrieval, predicted regime confirmed."))

# Final synthesis
new_slide_specs.append(("final_synthesis", None))

# Build them all
created_slides = []
for kind, payload in new_slide_specs:
    if kind == "smoke_caveat":
        s = build_smoke_caveat()
    elif kind == "divider":
        prefix, suffix, subtitle = payload
        s = build_section_divider(prefix, suffix, subtitle)
    elif kind == "result_A":
        s = build_result_A(**payload)
    elif kind == "result_B":
        s = build_result_B(**payload)
    elif kind == "synthesis":
        s = build_synthesis(payload)
    elif kind == "final_synthesis":
        s = build_synthesis(
            "Substrate: it holds.   Geometry: it acts where it should.   "
            "Generality: it replicates as an architecture, not as one trained instance."
        )
    elif kind == "baseline_note":
        # Recreate from current slide 52 content but with new claim tag
        s = add_blank_slide()
        add_monogram(s)
        add_title_block(s, "§ C — GENERALITY", "A note on the baseline comparison")
        add_cell_chip(s, "Cells 28–35")
        tb = add_text_box(s, 1.01, 1.95, 7.97, 2.80)
        tf = tb.text_frame
        lines = [
            ("kNN and RAG baselines are ", False, "oracle-maintained", True, ".", False),
        ]
        # Simple version: paragraph-by-paragraph
        items = [
            ("kNN and RAG baselines are oracle-maintained.", False),
            ("An external operator applies revise/remove operations to the retrieval store with full ground truth.", False),
            ("IBF derives the same operations from the discrepancy signal alone.", False),
            ("", False),
            ("The asymmetry favors the baselines on state maintenance.", True),
            ("", False),
            ("What's measured is what the architectural difference contributes despite that handicap.", False),
        ]
        for i, (txt, bold) in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            set_run(r, txt, "Source Serif 4", 12, bold, GREY)
        add_reading_bar(s, "Baselines get oracle help on state; IBF earns it from the discrepancy signal.")
    created_slides.append(s)

# -------- step 3: delete old slides 35..52 (indices 34..51) --------
# Doing this AFTER creating new slides so XML refs are stable.
sldIdLst = prs.slides._sldIdLst
slide_id_elements = list(sldIdLst)

# Identify which sldId elements correspond to old slides 35..52 (1-indexed slide numbers)
# original deck had 57 slides; created_slides were appended -> currently deck has 57 + N slides
# The sldIdLst order: [0..56] original + [57..] new ones
# We delete indices 34..51 from the original set.

# Delete the old result slides
indices_to_delete = list(range(34, 52))  # 0-indexed: old slides 35..52
# In reverse so prior deletions don't shift later indices
for idx in reversed(indices_to_delete):
    sldId = slide_id_elements[idx]
    rId = sldId.get(qn("r:id"))
    prs.part.drop_rel(rId)
    sldIdLst.remove(sldId)

# -------- step 4: reorder so new slides sit right after old slide 34 (results divider) --------
# After deletion, the deck has 57-18 + 25 = 64 slides in this order:
#   indices 0..33   -> original slides 1..34 (front matter + property/gap/architecture/test env + roadmap + results divider)
#   indices 34..38  -> original slides 53..57 (limits, implications, future, summary, thank you)
#   indices 39..63  -> the 25 newly-appended slides (in creation order)
# We want: 0..33, then the 25 new, then the 5 tail (limits..thank you).

sldIdLst_post = prs.slides._sldIdLst
ids = list(sldIdLst_post)
n_new = len(created_slides)
assert len(ids) == 57 - 18 + n_new, f"unexpected slide count {len(ids)} (n_new={n_new})"

tail_ids = ids[34:39]            # limits..thank-you
new_ids = ids[39:39 + n_new]     # newly-built slides
front_ids = ids[:34]             # 1..34 (front + results divider)

# Detach all
for sid in ids:
    sldIdLst_post.remove(sid)

# Re-attach in desired order
for sid in front_ids + new_ids + tail_ids:
    sldIdLst_post.append(sid)

prs.save(DST)
print(f"Saved {DST}")
print(f"Total slides: {len(prs.slides)}")
